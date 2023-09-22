import streamlit as st
from pptx import Presentation
from docx import Document
# import os
from pathlib import Path

uploaded_files = st.file_uploader("Choose a PPTX file",type=["pptx"],
                                  accept_multiple_files=True)

def upload():
    if not uploaded_files:
        st.text('Load pptx file/s first')
    else:
        for uploaded_file in uploaded_files:
            prs = Presentation(uploaded_file)
            st.header(uploaded_file.name)
            # Word
            document = Document()
            document.add_heading(uploaded_file.name, 0)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        st.write(shape.text)
                        cleaned_string = ''.join(c for c in shape.text if valid_xml_char_ordinal(c))
                        p = document.add_paragraph(cleaned_string)
            if save_doc_on:
                # document.save(uploaded_file.name + '.docx')
                
                # save_folder = '/Users/mshefi/downloads'
                if save_folder:
                    save_path = Path(save_folder, uploaded_file.name)
                    with open(save_path, mode='wb') as w:
                        w.write(uploaded_file.getvalue())

                    if save_path.exists():
                        st.success(f'File: {uploaded_file.name} is successfully saved!')
                else:
                    st.error('Input a folder to save')

def valid_xml_char_ordinal(c):
    codepoint = ord(c)
    # conditions ordered by presumed frequency
    return (
        0x20 <= codepoint <= 0xD7FF or
        codepoint in (0x9, 0xA, 0xD) or
        0xE000 <= codepoint <= 0xFFFD or
        0x10000 <= codepoint <= 0x10FFFF
        )

st.button("Parse text out of pptx file/s", on_click=upload, disabled=not uploaded_files)
save_doc_on = st.toggle('Save into a Word doc/s?')
if save_doc_on:
    save_folder = st.text_input('Save to folder')

# filename = file_selector()
# if filename:
#     st.write('You selected `%s`' % filename)
# else:
#     st.write('No file was selected')
# txt = st.text_area('Text to analyze', '''
#     It was the best of times, it was the worst of times, it was
#     the age of wisdom, it was the age of foolishness, it was
#     the epoch of belief, it was the epoch of incredulity, it
#     was the season of Light, it was the season of Darkness, it
#     was the spring of hope, it was the winter of despair, (...)
#     ''')